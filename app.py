import hashlib
import io
import json
import time

import numpy as np
import pandas as pd
import requests
import streamlit as st
from sentence_transformers import SentenceTransformer
from sklearn.cluster import KMeans
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

GROQ_CHAT_URL = "https://api.groq.com/openai/v1/chat/completions"
DEFAULT_GROQ_MODEL = "llama-3.3-70b-versatile"

OUTPUT_SCHEMA = {
    "type": "object",
    "properties": {
        "theme_name": {"type": "string"},
        "problem_summary": {"type": "string"},
        "sentiment": {"type": "string", "enum": ["positive", "neutral", "negative"]},
        "opportunity": {"type": "string"},
        "next_step": {"type": "string"},
        "success_metric": {
            "type": "string",
            "enum": ["Activation rate", "Time-to-value", "Error rate", "Retention", "CSAT", "Support ticket rate"],
        },
        "owner": {"type": "string", "enum": ["PM", "Eng", "Design", "Support"]},
        "confidence": {"type": "string", "enum": ["low", "medium", "high"]},
    },
    "required": [
        "theme_name",
        "problem_summary",
        "sentiment",
        "opportunity",
        "next_step",
        "success_metric",
        "owner",
        "confidence",
    ],
    "additionalProperties": False,
}

st.set_page_config(page_title="VoC Insight Agent", layout="wide")
st.title("VoC Insight Hub")

# Persist results across reruns (selectbox clicks, etc.) [web:193]
if "has_results" not in st.session_state:
    st.session_state["has_results"] = False
if "themes_df" not in st.session_state:
    st.session_state["themes_df"] = None
if "clustered_df" not in st.session_state:
    st.session_state["clustered_df"] = None
if "last_file_hash" not in st.session_state:
    st.session_state["last_file_hash"] = None


# ---------- Helpers ----------
@st.cache_resource
def load_embedder():
    return SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")


def file_bytes_and_hash(uploaded_file):
    uploaded_file.seek(0)
    data = uploaded_file.getvalue()
    if not data:
        raise ValueError("Uploaded file is empty (0 bytes). Please re-upload.")
    h = hashlib.sha256(data).hexdigest()
    return data, h


def read_uploaded_csv(uploaded_file) -> pd.DataFrame:
    data, _ = file_bytes_and_hash(uploaded_file)
    return pd.read_csv(io.BytesIO(data))


@st.cache_data(show_spinner=False)
def cached_embeddings(text_list: list[str], file_hash: str, max_rows: int) -> np.ndarray:
    embedder = load_embedder()
    emb = embedder.encode(text_list, show_progress_bar=False, normalize_embeddings=True)
    return np.array(emb)


@st.cache_data(ttl=6 * 60 * 60, show_spinner=False)
def groq_generate_json(prompt: str, schema: dict, model: str) -> dict:
    api_key = st.secrets["GROQ_API_KEY"]

    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": "Return ONLY valid JSON that matches the provided schema."},
            {"role": "user", "content": prompt},
        ],
        "temperature": 0,
        "response_format": {"type": "json_object"},
    }

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    r = requests.post(GROQ_CHAT_URL, headers=headers, json=payload, timeout=120)
    if not r.ok:
        raise RuntimeError(f"Groq HTTP {r.status_code}: {r.text}")
    data = r.json()

    content = data["choices"][0]["message"]["content"]
    return json.loads(content)


@st.cache_data
def df_to_csv_bytes(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=False).encode("utf-8")


def severity_to_score(x):
    if pd.isna(x):
        return np.nan
    s = str(x).strip().lower()
    return {"low": 1, "medium": 2, "high": 3, "critical": 4}.get(s, np.nan)


def plan_to_weight(x):
    if pd.isna(x):
        return 1.0
    s = str(x).strip().lower()
    return {"free": 1.0, "pro": 1.5, "team": 2.0, "enterprise": 2.5}.get(s, 1.0)


def prettify_columns(df_in: pd.DataFrame) -> pd.DataFrame:
    df2 = df_in.copy()
    df2.columns = [c.replace("_", " ").strip().title() for c in df2.columns]
    return df2


def make_column_config(pretty_df: pd.DataFrame) -> dict:
    # Column config controls label + relative width [web:127]
    cols = set(pretty_df.columns)
    cfg = {}

    def add(name: str, label: str, width: str):
        if name in cols:
            cfg[name] = st.column_config.Column(label=label, width=width)

    # Smaller columns
    add("Theme Name", "Theme Name", "medium")
    add("Priority Score", "Priority Score", "small")
    add("Count", "Count", "small")
    add("Owner", "Owner", "small")
    add("Success Metric", "Success Metric", "medium")
    add("Confidence", "Confidence", "small")
    add("Sentiment", "Sentiment", "small")
    add("Cluster", "Cluster", "small")

    add("Feedback Id", "Feedback ID", "medium")
    add("Product", "Product", "medium")
    add("Date", "Date", "small")
    add("Source", "Source", "medium")
    add("Persona", "Persona", "small")
    add("Plan", "Plan", "small")
    add("Severity", "Severity", "small")

    # Wider text columns
    add("Problem Summary", "Problem Summary", "large")
    add("Opportunity", "Opportunity", "large")
    add("Next Step", "Next Step", "large")
    add("Evidence Quotes", "Evidence Quotes", "large")
    add("Text", "Text", "large")

    add("Avg Severity", "Avg Severity", "small")
    add("Avg Plan Weight", "Avg Plan Weight", "small")

    return cfg


def create_ppt_summary(themes_df: pd.DataFrame) -> bytes:
    """Generate PPT with title slide + top 5 themes summary + detail slides [web:234][web:239]"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "VoC Insight Summary"
    subtitle.text = f"Generated {time.strftime('%B %d, %Y')}"

    # Slide 2: Top 5 themes table
    bullet_slide_layout = prs.slide_layouts[5]  # Title only layout
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Top 5 Themes by Priority"

    # Add table [web:247]
    rows = 6  # header + 5 themes
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths [web:289]
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(1.2)
    table.columns[2].width = Inches(1)
    table.columns[3].width = Inches(4.3)  # Make "Next Step" column wider

    # Header row
    headers = ["Theme", "Priority", "Owner", "Next Step"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Data rows
    top5 = themes_df.head(5)
    for row_idx, (_, row) in enumerate(top5.iterrows(), start=1):
        # Theme name (no truncation)
        cell = table.cell(row_idx, 0)
        cell.text = str(row["theme_name"])
        cell.text_frame.word_wrap = True  # Enable wrapping [web:281]
        cell.text_frame.paragraphs[0].font.size = Pt(10)

        # Priority
        cell = table.cell(row_idx, 1)
        cell.text = f"{row['priority_score']:.1f}"
        cell.text_frame.paragraphs[0].font.size = Pt(10)

        # Owner
        cell = table.cell(row_idx, 2)
        cell.text = str(row["owner"])
        cell.text_frame.paragraphs[0].font.size = Pt(10)

        # Next Step (FULL TEXT - no truncation) [web:281]
        cell = table.cell(row_idx, 3)
        cell.text = str(row["next_step"])  # No [:60] truncation
        cell.text_frame.word_wrap = True  # Enable wrapping
        cell.text_frame.paragraphs[0].font.size = Pt(9)  # Slightly smaller font

    # Slides 3+: One slide per top 5 theme [web:238]
    for _, row in top5.iterrows():
        bullet_slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(bullet_slide_layout)

        title = slide.shapes.title
        title.text = str(row["theme_name"])

        # Add bullet points
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        # Problem
        p = tf.paragraphs[0]
        p.text = f"Problem: {row['problem_summary']}"
        p.level = 0
        p.font.size = Pt(14)

        # Opportunity
        p = tf.add_paragraph()
        p.text = f"{row['opportunity']}"
        p.level = 0
        p.font.size = Pt(14)

        # Next step (FULL TEXT)
        p = tf.add_paragraph()
        p.text = f"Next Step: {row['next_step']}"  # No truncation
        p.level = 0
        p.font.size = Pt(14)

        # Metrics
        p = tf.add_paragraph()
        p.text = f"Success Metric: {row['success_metric']} | Owner: {row['owner']} | Confidence: {row['confidence']}"
        p.level = 1
        p.font.size = Pt(12)

    # Save to BytesIO [web:239]
    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output.getvalue()



# ---------- UI controls ----------
from pathlib import Path

APP_DIR = Path(__file__).resolve().parent
SAMPLE_PATH = APP_DIR / "sample_feedback.csv"

if SAMPLE_PATH.exists():
    st.download_button(
        "Download sample_feedback.csv",
        data=SAMPLE_PATH.read_bytes(),
        file_name="sample_feedback.csv",
        mime="text/csv",
        use_container_width=True,
    )
    st.caption("Tip: Download this file, then upload it to try the app quickly.")

uploaded = st.file_uploader("Upload feedback CSV (needs a 'text' column)", type=["csv"])

k = st.slider("Number of themes", 2, 12, 8)
max_rows = st.slider("Max rows to analyze (speed)", 50, 2000, 200)

if uploaded is None:
    st.info("Upload sample_feedback.csv")
    st.stop()

st.write("Uploaded file:", uploaded.name)

try:
    df_raw = read_uploaded_csv(uploaded)
    _, file_hash = file_bytes_and_hash(uploaded)
except Exception as e:
    st.error(str(e))
    st.stop()

if "text" not in df_raw.columns:
    st.error("CSV must contain a 'text' column.")
    st.stop()

# If user uploads a new file, invalidate prior results automatically
if st.session_state["last_file_hash"] is None:
    st.session_state["last_file_hash"] = file_hash
elif st.session_state["last_file_hash"] != file_hash:
    st.session_state["has_results"] = False
    st.session_state["themes_df"] = None
    st.session_state["clustered_df"] = None
    st.session_state["last_file_hash"] = file_hash

df = df_raw.copy()
with st.expander("Filters", expanded=True):
    for colname in ["product", "persona", "plan", "source", "severity"]:
        if colname in df.columns:
            options = sorted([x for x in df[colname].dropna().unique()])
            chosen = st.multiselect(colname.capitalize(), options)
            if chosen:
                df = df[df[colname].isin(chosen)]

df = df.dropna(subset=["text"]).copy()
df["text"] = df["text"].astype(str)
df = df.head(max_rows)

if len(df) == 0:
    st.warning("No rows left after filtering. Clear filters or upload more data.")
    st.stop()

st.subheader("Preview")
preview_pretty = prettify_columns(df.head(20))
st.dataframe(
    preview_pretty,
    use_container_width=True,
    column_config=make_column_config(preview_pretty),
    hide_index=True,
)

# Buttons below preview [web:219]
col_a, col_b = st.columns([1, 1])
with col_a:
    if st.button("Generate insights", type="primary", use_container_width=True):
        st.session_state["has_results"] = True
with col_b:
    if st.button("Reset / Run again", use_container_width=True):
        st.session_state["has_results"] = False
        st.session_state["themes_df"] = None
        st.session_state["clustered_df"] = None
        st.session_state["last_file_hash"] = None
        st.rerun()

# If we already have results, just render them (don't force rerun)
if st.session_state["has_results"] and st.session_state["themes_df"] is not None and st.session_state["clustered_df"] is not None:
    themes_df = st.session_state["themes_df"]
    df_clustered = st.session_state["clustered_df"]
else:
    if not st.session_state["has_results"]:
        st.info("Click \"Generate insights\" to run clustering + Groq labeling.")
        st.stop()

    # ---------- Debug checkpoints ----------
    st.info("Step 1/3: Starting embedding + clustering...")

    # ---------- Embedding + clustering ----------
    st.info("Step 2/3: Embedding texts now (first run may download model)...")
    with st.spinner("Embedding + clustering..."):
        texts = df["text"].tolist()
        emb = cached_embeddings(texts, file_hash, max_rows)

        n_samples = emb.shape[0]
        if n_samples < 2:
            st.error("Not enough rows to cluster after filtering. Remove filters or increase max rows.")
            st.stop()

        k_eff = min(k, n_samples)
        if k_eff != k:
            st.warning(f"Reducing number of themes from {k} to {k_eff} because only {n_samples} rows are available.")

        km = KMeans(n_clusters=k_eff, n_init="auto", random_state=42)
        df = df.copy()
        df["cluster"] = km.fit_predict(emb)

    model = st.secrets.get("GROQ_MODEL", DEFAULT_GROQ_MODEL)
    st.caption(f"Run settings: k={k_eff}, rows={len(df)}, model={model}")

    # ---------- Theme scoring + labeling ----------
    st.info("Step 3/3: Starting Groq labeling...")

    themes = []
    progress = st.progress(0)
    status = st.empty()

    clusters = sorted(df["cluster"].unique())
    total = len(clusters)

    for idx, c in enumerate(clusters, start=1):
        dfc = df[df["cluster"] == c].copy()

        quotes = [q[:400] for q in dfc["text"].head(5).tolist()]
        evidence = " | ".join([q[:140].replace("\\n", " ") for q in quotes])

        count = int(len(dfc))

        if "severity" in dfc.columns:
            dfc["_severity_score"] = dfc["severity"].apply(severity_to_score)
            avg_sev = float(np.nanmean(dfc["_severity_score"])) if dfc["_severity_score"].notna().any() else np.nan
        else:
            avg_sev = np.nan

        if "plan" in dfc.columns:
            dfc["_plan_weight"] = dfc["plan"].apply(plan_to_weight)
            avg_plan_weight = float(np.nanmean(dfc["_plan_weight"])) if dfc["_plan_weight"].notna().any() else 1.0
        else:
            avg_plan_weight = 1.0

        sev_component = (avg_sev if not np.isnan(avg_sev) else 1.5)
        priority_score = count * sev_component * avg_plan_weight

        prompt = f"""
You are a Product Manager analyzing customer feedback.

Return ONLY a JSON object with exactly these keys:
theme_name, problem_summary, sentiment, opportunity, next_step, success_metric, owner, confidence.

Rules:
- theme_name: 2-5 words, noun phrase, avoid "issues"
- problem_summary: exactly 2 sentences, do NOT copy snippets verbatim, no quotes
- sentiment: one of ["positive","neutral","negative"]
- opportunity: one sentence starting with "Opportunity:"
- next_step: one sentence, must start with an action verb (e.g., "Instrument", "Interview", "Reproduce", "A/B test", "Audit")
- owner: choose ONLY one from ["PM","Eng","Design","Support"] based on who would do the next_step
- success_metric: choose ONLY one from ["Activation rate","Time-to-value","Error rate","Retention","CSAT","Support ticket rate"]
- confidence: choose ONLY one from ["low","medium","high"] using this rubric:
  - high = snippets are specific and consistent (same problem repeated)
  - medium = theme is clear but details vary
  - low = ambiguous or mixed topics
Return ONLY JSON (no markdown, no extra text).

Snippets:
{json.dumps(quotes, ensure_ascii=False)}
""".strip()

        status.info(f"Calling Groq for theme {idx}/{total} (cluster {c})...")
        time.sleep(0.4)

        try:
            model = st.secrets.get("GROQ_MODEL", DEFAULT_GROQ_MODEL)
            obj = groq_generate_json(prompt, OUTPUT_SCHEMA, model)
        except Exception as e:
            st.error(f"Groq call failed for cluster {c}: {e}")
            obj = {
                "theme_name": f"Theme {c}",
                "problem_summary": "Groq call failed. Theme label is a placeholder.",
                "sentiment": "neutral",
                "opportunity": "Opportunity: Retry with fewer rows/themes or increase timeout.",
                "next_step": "Retry labeling after checking Groq / rate limits.",
                "success_metric": "Activation rate",
                "owner": "PM",
                "confidence": "low",
            }

        themes.append(
            {
                "cluster": int(c),
                "count": count,
                "avg_severity": avg_sev,
                "avg_plan_weight": avg_plan_weight,
                "priority_score": float(priority_score),
                "theme_name": obj.get("theme_name"),
                "problem_summary": obj.get("problem_summary"),
                "sentiment": obj.get("sentiment"),
                "opportunity": obj.get("opportunity"),
                "next_step": obj.get("next_step"),
                "success_metric": obj.get("success_metric"),
                "owner": obj.get("owner"),
                "confidence": obj.get("confidence"),
                "evidence_quotes": evidence,
            }
        )

        progress.progress(int(idx / total * 100))

    status.success("Groq labeling complete.")

    themes_df = (
        pd.DataFrame(themes)
        .sort_values("priority_score", ascending=False)
        .reset_index(drop=True)
    )

    st.session_state["themes_df"] = themes_df
    st.session_state["clustered_df"] = df

    df_clustered = df

# ---------- Outputs ----------
st.subheader("Top 5 themes (by priority)")
top5 = themes_df[
    ["theme_name", "priority_score", "count", "owner", "success_metric", "confidence", "opportunity"]
].head(5)
top5_pretty = prettify_columns(top5)
st.dataframe(
    top5_pretty,
    use_container_width=True,
    column_config=make_column_config(top5_pretty),
    hide_index=True,
)

st.subheader("Themes (full)")
full_cols = [
    "theme_name",
    "priority_score",
    "count",
    "owner",
    "success_metric",
    "confidence",
    "sentiment",
    "problem_summary",
    "opportunity",
    "next_step",
    "avg_severity",
    "avg_plan_weight",
    "evidence_quotes",
]
full_df = themes_df[full_cols]
full_pretty = prettify_columns(full_df)
st.dataframe(
    full_pretty,
    use_container_width=True,
    column_config=make_column_config(full_pretty),
    hide_index=True,
)

st.subheader("Inspect a theme (drill-down)")
choice = st.selectbox("Theme", themes_df["theme_name"].tolist())
sel_cluster = int(themes_df.loc[themes_df["theme_name"] == choice, "cluster"].iloc[0])

cols_to_show = [
    c for c in ["feedback_id", "product", "date", "source", "persona", "plan", "severity", "text"] if c in df_clustered.columns
]
if not cols_to_show:
    cols_to_show = ["text"]

drill_df = df_clustered[df_clustered["cluster"] == sel_cluster][cols_to_show].head(50)
drill_pretty = prettify_columns(drill_df)
st.dataframe(
    drill_pretty,
    use_container_width=True,
    column_config=make_column_config(drill_pretty),
    hide_index=True,
)

# Download buttons
st.subheader("Export")
col1, col2, col3 = st.columns(3)

with col1:
    st.download_button(
        "📊 Download themes CSV",
        data=df_to_csv_bytes(full_df),
        file_name="themes.csv",
        mime="text/csv",
        use_container_width=True,
    )

with col2:
    st.download_button(
        "📋 Download clustered feedback CSV",
        data=df_to_csv_bytes(df_clustered),
        file_name="clustered_feedback.csv",
        mime="text/csv",
        use_container_width=True,
    )

with col3:
    ppt_bytes = create_ppt_summary(themes_df)
    st.download_button(
        "📽️ Download PPT summary",
        data=ppt_bytes,
        file_name="voc_insights_summary.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
    )



